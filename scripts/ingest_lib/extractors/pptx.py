"""PPTX extractor: one Markdown section per slide."""
from __future__ import annotations

from pathlib import Path

from .base import ExtractionResult


def extract(src: Path, _assets_dir: Path) -> ExtractionResult:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:
        return ExtractionResult(
            status="manual_review",
            extractor="pptx",
            markdown="",
            error=f"python-pptx missing: {exc}",
        )

    try:
        prs = Presentation(str(src))
    except Exception as exc:  # noqa: BLE001
        return ExtractionResult(
            status="manual_review",
            extractor="pptx",
            markdown="",
            error=f"open failed: {exc}",
        )

    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide) or f"Slide {i}"
        parts.append(f"## {i}. {title}")
        body_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs).strip()
                if txt:
                    body_lines.append(txt)
        if body_lines:
            parts.append("\n".join(f"- {ln}" for ln in body_lines))
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_txt = slide.notes_slide.notes_text_frame.text.strip()
            if notes_txt:
                parts.append(f"**Notes:** {notes_txt}")
        parts.append("")
    return ExtractionResult(
        status="processed",
        extractor="pptx",
        markdown="\n".join(parts),
    )


def _slide_title(slide) -> str | None:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip() or None
    return None
